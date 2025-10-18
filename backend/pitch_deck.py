"""
Pitch Deck Generator using Perplexity Sonar Pro + python-pptx
Sonar Pro generates the content, python-pptx creates the presentation
"""

import os
import json
import asyncio
from datetime import datetime
from dotenv import load_dotenv
from perplexity import AsyncPerplexity

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  Warning: python-pptx not installed. Install with: pip install python-pptx")

load_dotenv()

DEFAULT_STARTUP_IDEA = "AI-powered legal technology platform for contract review"

async def generate_pitch_content(client, startup_idea: str) -> dict:
    """Generate comprehensive pitch deck content using Perplexity Sonar Pro"""
    
    prompt = f"""
    Create a comprehensive pitch deck for the following startup idea:
    "{startup_idea}"
    
    Generate detailed content for each of these 12 slides:
    
    1. TITLE SLIDE
       - Company name (create a compelling name based on the idea)
       - Tagline (one compelling sentence)
       - Founder placeholder text
    
    2. PROBLEM
       - What problem does this solve?
       - Include 3-4 specific pain points
       - Add real statistics if available
    
    3. SOLUTION
       - How does this startup solve the problem?
       - Key features (3-4 bullet points)
       - Unique approach
    
    4. MARKET SIZE
       - Total Addressable Market (TAM)
       - Serviceable Addressable Market (SAM)
       - Serviceable Obtainable Market (SOM)
       - Include dollar amounts and growth rates
    
    5. PRODUCT
       - Product/service description
       - Key features (4-5 points)
       - Technology stack or approach
    
    6. BUSINESS MODEL
       - Revenue streams
       - Pricing strategy
       - Unit economics
       - Path to profitability
    
    7. TRACTION
       - Key milestones achieved or planned
       - Early wins or validation
       - Customer testimonials or LOIs
    
    8. COMPETITION
       - Main competitors (3-4)
       - Competitive advantages (3-4 points)
       - What makes this different
    
    9. GO-TO-MARKET STRATEGY
       - Customer acquisition channels
       - Marketing approach
       - Sales strategy
       - Partnerships
    
    10. TEAM
        - Required expertise
        - Key roles needed
        - Advisory board possibilities
    
    11. FINANCIAL PROJECTIONS
        - 3-year revenue forecast
        - Key metrics
        - Growth assumptions
    
    12. FUNDING ASK
        - Amount seeking
        - Use of funds (breakdown)
        - Milestones to achieve
        - Expected outcomes
    
    Return ONLY a valid JSON object with this EXACT structure:
    {{
        "company_name": "Company Name Here",
        "tagline": "One-line tagline",
        "slides": [
            {{
                "title": "Problem",
                "bullets": [
                    "First key point",
                    "Second key point",
                    "Third key point"
                ]
            }},
            {{
                "title": "Solution",
                "bullets": [
                    "Key solution point 1",
                    "Key solution point 2"
                ]
            }}
        ]
    }}
    
    IMPORTANT: Each slide should have 3-5 concise bullet points (not paragraphs).
    Make content compelling, data-driven, and investor-ready.
    """
    
    response = await client.chat.completions.create(
        model="sonar-pro",
        messages=[
            {
                "role": "system",
                "content": """You are an expert startup advisor and pitch deck creator. 
You create compelling, concise pitch deck content that investors love.
Return ONLY valid JSON with no additional text before or after.
Each slide should have clear, concise bullet points - not paragraphs."""
            },
            {
                "role": "user",
                "content": prompt
            }
        ]
    )
    
    return response.choices[0].message.content

def extract_json_from_response(text: str) -> dict:
    """Extract JSON object from Perplexity's response"""
    import re
    
    # Try to find JSON object in the response
    json_match = re.search(r'\{.*\}', text, re.DOTALL)
    if json_match:
        try:
            data = json.loads(json_match.group())
            return data
        except json.JSONDecodeError as e:
            print(f"   ⚠️  JSON parse error: {e}")
            print(f"   Response preview: {text[:500]}...")
            return None
    
    print(f"   ⚠️  No JSON object found in response")
    return None

def create_presentation(pitch_data: dict, output_file: str):
    """Create PowerPoint presentation using python-pptx"""
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx not available. Cannot create presentation.")
        return False
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    company_name = pitch_data.get('company_name', 'Startup Name')
    tagline = pitch_data.get('tagline', '')
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = company_name
    subtitle.text = tagline
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # Add content slides
    slides_data = pitch_data.get('slides', [])
    
    for slide_data in slides_data:
        slide_title = slide_data.get('title', 'Untitled')
        bullets = slide_data.get('bullets', [])
        
        # Use title and content layout
        bullet_slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add bullets to content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()  # Clear default text
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(12)
    
    # Add closing slide
    closing_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(closing_slide_layout)
    
    # Add centered text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    p2 = text_frame.add_paragraph()
    p2.text = "Questions?"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(36)
    p2.font.color.rgb = RGBColor(102, 102, 102)
    
    # Save presentation
    prs.save(output_file)
    print(f"✅ PowerPoint saved to: {output_file}")
    return True

async def generate_pitch_deck(startup_idea: str) -> dict:
    """
    Generate a complete pitch deck for a startup idea
    
    Args:
        startup_idea: Description of the startup idea
    
    Returns:
        dict with pitch content and file paths
    """
    
    print(f"\n{'='*80}")
    print(f"🎯 GENERATING PITCH DECK FOR:")
    print(f"{startup_idea}")
    print(f"{'='*80}\n")
    
    # Step 1: Generate content using Perplexity Sonar Pro
    print("📝 Step 1: Generating pitch deck content with Perplexity Sonar Pro...")
    
    async with AsyncPerplexity() as client:
        response = await generate_pitch_content(client, startup_idea)
    
    print("✅ Content generated!\n")
    
    # Step 2: Parse the response
    print("📊 Step 2: Parsing content...")
    pitch_data = extract_json_from_response(response)
    
    if not pitch_data:
        print("❌ Failed to parse pitch deck content")
        return {
            "success": False,
            "error": "Failed to parse content",
            "raw_response": response
        }
    
    print("✅ Content parsed!\n")
    
    # Step 3: Save JSON version
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    json_filename = f"pitch_deck_{timestamp}.json"
    
    with open(json_filename, "w", encoding="utf-8") as f:
        json.dump(pitch_data, f, indent=2, ensure_ascii=False)
    
    print(f"💾 Saved pitch data to: {json_filename}\n")
    
    # Step 4: Generate PowerPoint presentation
    pptx_filename = None
    
    if PPTX_AVAILABLE:
        print("🎨 Step 3: Generating PowerPoint presentation with python-pptx...")
        
        pptx_filename = f"pitch_deck_{timestamp}.pptx"
        success = create_presentation(pitch_data, pptx_filename)
        
        if success:
            print(f"✅ Presentation created successfully!\n")
        else:
            print(f"⚠️  Failed to create presentation\n")
    else:
        print("ℹ️  python-pptx not available. Install with: pip install python-pptx")
        print("   Content has been saved as JSON file.\n")
    
    # Step 5: Display summary
    print(f"{'='*80}")
    print(f"📊 PITCH DECK SUMMARY")
    print(f"{'='*80}")
    print(f"Company: {pitch_data.get('company_name', 'N/A')}")
    print(f"Tagline: {pitch_data.get('tagline', 'N/A')}")
    print(f"Number of slides: {len(pitch_data.get('slides', [])) + 2}")  # +2 for title and closing
    print(f"\n📁 Generated files:")
    print(f"   • {json_filename} (JSON data)")
    if pptx_filename:
        print(f"   • {pptx_filename} (PowerPoint)")
    print(f"{'='*80}\n")
    
    # Display slide titles
    slides = pitch_data.get('slides', [])
    if slides:
        print("📑 SLIDE OUTLINE:\n")
        print("1. Title Slide")
        for i, slide in enumerate(slides, 2):
            print(f"{i}. {slide.get('title', 'Untitled')}")
        print(f"{len(slides) + 2}. Thank You")
        print()
    
    # Display content preview
    if slides:
        print(f"{'='*80}")
        print("📄 CONTENT PREVIEW - First 3 Slides")
        print(f"{'='*80}\n")
        
        for i, slide in enumerate(slides[:3], 1):
            print(f"SLIDE {i + 1}: {slide.get('title', 'Untitled')}")
            print("-" * 80)
            bullets = slide.get('bullets', [])
            for bullet in bullets:
                print(f"  • {bullet}")
            print()
    
    return {
        "success": True,
        "pitch_data": pitch_data,
        "json_file": json_filename,
        "pptx_file": pptx_filename
    }

async def generate_pitch_deck_api(startup_idea: str) -> dict:
    """
    API-friendly version of pitch deck generator
    Returns structured data ready for JSON serialization
    """
    result = await generate_pitch_deck(startup_idea)
    return result

async def main():
    """Interactive CLI for pitch deck generation"""
    
    print("\n" + "="*80)
    print("🎯 PITCH DECK GENERATOR")
    print("="*80)
    print("This tool generates a complete investor pitch deck using:")
    print("  • Perplexity Sonar Pro (content generation)")
    print("  • python-pptx (PowerPoint creation)")
    print()
    
    # Check for python-pptx
    if not PPTX_AVAILABLE:
        print("⚠️  python-pptx is not installed!")
        print("   Install it with: pip install python-pptx")
        print("   (You can still generate JSON content without it)\n")
    
    # Get startup idea from user
    startup_idea = input(f"Enter your startup idea\n(default: '{DEFAULT_STARTUP_IDEA}'): ").strip()
    if not startup_idea:
        startup_idea = DEFAULT_STARTUP_IDEA
    
    # Generate the pitch deck
    result = await generate_pitch_deck(startup_idea)
    
    if result.get('success'):
        print("✅ Pitch deck generated successfully!")
        
        if result.get('pptx_file'):
            print(f"\n🎉 Open your presentation: {result.get('pptx_file')}")
        else:
            print(f"\n📄 View the data: {result.get('json_file')}")
    else:
        print("❌ Failed to generate pitch deck")
        print(f"Error: {result.get('error', 'Unknown error')}")

if __name__ == "__main__":
    # Check for required environment variables
    if not os.environ.get("PERPLEXITY_API_KEY"):
        print("⚠️  Warning: PERPLEXITY_API_KEY not found in environment variables")
        print("   Please set it in your .env file or environment")
        print()
    
    asyncio.run(main())
